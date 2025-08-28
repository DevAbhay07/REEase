import streamlit as st
import PyPDF2
import json
import os
from dotenv import load_dotenv
import io
from summarizer import TextSummarizer
import docx
import pptx
load_dotenv()

# Handle textract import gracefully
try:
    import textract
    HAS_TEXTRACT = True
except ImportError:
    HAS_TEXTRACT = False

def extract_text_from_file(uploaded_file) -> str:
    """Extract text from various file formats"""
    text = ''
    try:
        file_extension = uploaded_file.name.split('.')[-1].lower()
        
        if file_extension == "pdf":
            pdf_reader = PyPDF2.PdfReader(uploaded_file)
            for page_num in range(len(pdf_reader.pages)):
                page_text = pdf_reader.pages[page_num].extract_text()
                if page_text:
                    text += page_text + "\n"
                    
        elif file_extension == "docx":
            doc = docx.Document(uploaded_file)
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text += paragraph.text + "\n"
                    
        elif file_extension == "txt":
            text = uploaded_file.read().decode("utf-8")
            
        elif file_extension == "pptx":
            ppt = pptx.Presentation(uploaded_file)
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text += shape.text + "\n"
        else:
            # Fallback to textract for other formats
            if HAS_TEXTRACT:
                try:
                    text = textract.process(uploaded_file).decode("utf-8")
                except Exception:
                    st.error(f"Error processing {file_extension.upper()} file with textract")
                    return ""
            else:
                st.error(f"File format {file_extension.upper()} not supported. Please use PDF, DOCX, TXT, PPTX, or JSON files.")
                return ""
            
    except Exception as e:
        st.error(f"Error extracting text from {uploaded_file.name}: {str(e)}")
        return ""
    
    return text.strip()

def process_json_file(uploaded_file) -> list:
    """Process uploaded JSON file and extract thread data"""
    try:
        # Read the JSON file content
        content = uploaded_file.read().decode("utf-8")
        data = json.loads(content)
        
        if not isinstance(data, list):
            st.error("JSON file should contain a list of email messages")
            return []
        
        return data
        
    except json.JSONDecodeError as e:
        st.error(f"Invalid JSON file: {str(e)}")
        return []
    except Exception as e:
        st.error(f"Error processing JSON file: {str(e)}")
        return []

def create_download_file(content: str, file_format: str) -> bytes:
    """Create downloadable file content"""
    try:
        if file_format.lower() == "docx":
            doc = docx.Document()
            doc.add_paragraph(content)
            
            # Save to bytes buffer
            buffer = io.BytesIO()
            doc.save(buffer)
            buffer.seek(0)
            return buffer.getvalue()
        else:
            # Default to text format
            return content.encode('utf-8')
    except Exception as e:
        st.error(f"Error creating download file: {str(e)}")
        return content.encode('utf-8')

def main():
    # Page config
    st.set_page_config(
        page_title="REEASE - AI Document Summarizer",
        page_icon="🤖",
        layout="wide",
        initial_sidebar_state="expanded"
    )
    
    # Sidebar
    with st.sidebar:
        st.markdown("# 🤖 REEASE")
        st.markdown("### AI Document Summarizer")
        st.markdown("---")
        
        st.markdown("### 📊 Features")
        st.markdown("""
        ✅ **Multi-Format Support**  
        📄 PDF, DOCX, TXT, PPTX, JSON
        
        ✅ **Smart Compression**  
        🎯 25%, 50%, 75% precision control
        
        ✅ **Email Intelligence**  
        📧 Thread-aware processing
        
        ✅ **Lightning Fast**  
        ⚡ 3-5 second processing
        """)
        
        st.markdown("---")
        st.markdown("### 🌟 Stats")
        
        col1, col2 = st.columns(2)
        with col1:
            st.metric("Formats", "5+", "PDF, DOCX...")
        with col2:
            st.metric("Speed", "3-5s", "⚡ Fast")
            
        col3, col4 = st.columns(2)
        with col3:
            st.metric("AI Model", "Gemini", "2.5 Flash")
        with col4:
            st.metric("Platforms", "2", "Web + Extension")
        
        st.markdown("---")
        st.markdown("### 🔥 What's New")
        st.info("🚀 Chrome Extension now available with browser integration!")
        
        st.markdown("---")
        st.markdown("### 💡 Pro Tips")
        st.markdown("""
        🎯 **Use 25%** for quick overviews  
        📝 **Use 50%** for balanced summaries  
        📚 **Use 75%** for detailed insights  
        📧 **JSON mode** for email threads
        """)
    
    # Main header with better styling
    st.markdown("""
    <div style='text-align: center; padding: 2rem 0;'>
        <h1 style='color: #1f77b4; font-size: 3rem; margin-bottom: 0.5rem;'>
            🤖 REEASE
        </h1>
        <h3 style='color: #666; font-weight: 300; margin-bottom: 1rem;'>
            AI-Powered Document Summarization Platform
        </h3>
        <p style='color: #888; font-size: 1.1rem; max-width: 600px; margin: 0 auto;'>
            Transform any document into intelligent summaries with precision control. 
            Powered by Google Gemini AI for lightning-fast results.
        </p>
    </div>
    """, unsafe_allow_html=True)
    
    # Feature highlights
    col1, col2, col3, col4 = st.columns(4)
    
    with col1:
        st.markdown("""
        <div style='text-align: center; padding: 1rem; background: linear-gradient(45deg, #ff6b6b, #ee5a52); 
                    border-radius: 10px; color: white; margin-bottom: 1rem;'>
            <h3 style='margin: 0; color: white;'>📄</h3>
            <p style='margin: 0; color: white; font-weight: bold;'>Multi-Format</p>
            <small style='color: rgba(255,255,255,0.8);'>PDF, DOCX, TXT, PPTX</small>
        </div>
        """, unsafe_allow_html=True)
    
    with col2:
        st.markdown("""
        <div style='text-align: center; padding: 1rem; background: linear-gradient(45deg, #4ecdc4, #44a08d); 
                    border-radius: 10px; color: white; margin-bottom: 1rem;'>
            <h3 style='margin: 0; color: white;'>⚡</h3>
            <p style='margin: 0; color: white; font-weight: bold;'>Lightning Fast</p>
            <small style='color: rgba(255,255,255,0.8);'>3-5 seconds</small>
        </div>
        """, unsafe_allow_html=True)
    
    with col3:
        st.markdown("""
        <div style='text-align: center; padding: 1rem; background: linear-gradient(45deg, #667eea, #764ba2); 
                    border-radius: 10px; color: white; margin-bottom: 1rem;'>
            <h3 style='margin: 0; color: white;'>🎯</h3>
            <p style='margin: 0; color: white; font-weight: bold;'>Precision Control</p>
            <small style='color: rgba(255,255,255,0.8);'>25%, 50%, 75%</small>
        </div>
        """, unsafe_allow_html=True)
    
    with col4:
        st.markdown("""
        <div style='text-align: center; padding: 1rem; background: linear-gradient(45deg, #ffeaa7, #fdcb6e); 
                    border-radius: 10px; color: white; margin-bottom: 1rem;'>
            <h3 style='margin: 0; color: white;'>🤖</h3>
            <p style='margin: 0; color: white; font-weight: bold;'>Gemini AI</p>
            <small style='color: rgba(255,255,255,0.8);'>Latest Model</small>
        </div>
        """, unsafe_allow_html=True)
    
    st.markdown("---")
    
    # Check for API key
    if not os.environ.get("GEMINI_API_KEY"):
        st.error("⚠️ GEMINI_API_KEY environment variable is not set. Please configure your API key to use this application.")
        st.stop()
    
    try:
        summarizer = TextSummarizer()
    except Exception as e:
        st.error(f"Failed to initialize summarizer: {str(e)}")
        st.stop()
    
    # Main interface with better styling
    st.markdown("### 🚀 Choose Your Input Method")
    input_method = st.radio(
        "", 
        ("📁 Upload File", "✏️ Enter Text", "📧 Upload JSON (Email Threads)"),
        help="Choose how you want to provide content for summarization",
        horizontal=True
    )
    
    if input_method == "📁 Upload File":
        st.markdown("### 📁 File Upload")
        st.markdown("Drag and drop or browse to upload your document")
        
        uploaded_file = st.file_uploader(
            "Choose a file", 
            type=["pdf", "docx", "txt", "pptx"],
            help="Supported formats: PDF, Word, Text, PowerPoint"
        )
        
        if uploaded_file is not None:
            st.success(f"✅ File '{uploaded_file.name}' uploaded successfully!")
            
            # Summarization options
            col1, col2 = st.columns(2)
            with col1:
                compression = st.selectbox(
                    "Summary Length", 
                    ["Regular", "25%", "50%", "75%"],
                    help="Choose the compression level for your summary"
                )
            
            with col2:
                download_format = st.selectbox(
                    "Download Format", 
                    ["TXT", "DOCX"],
                    help="Choose the format for the downloaded summary"
                )
            
            if st.button("🚀 Generate Summary", type="primary"):
                with st.spinner("Extracting text from file..."):
                    extracted_text = extract_text_from_file(uploaded_file)
                
                if extracted_text:
                    st.success("✅ Text extracted successfully!")
                    
                    with st.spinner("Generating AI summary..."):
                        compression_ratio = None if compression == "Regular" else compression
                        summary = summarizer.generate_summary(extracted_text, compression_ratio)
                    
                    if summary:
                        st.markdown("### 📝 Generated Summary")
                        st.success("✨ **AI Summary Generated Successfully!**")
                        st.text_area("", summary, height=300, disabled=True, label_visibility="collapsed")
                        
                        # Download functionality
                        file_content = create_download_file(summary, download_format)
                        file_name = f"summary.{download_format.lower()}"
                        
                        st.download_button(
                            label=f"⬇️ Download {download_format}",
                            data=file_content,
                            file_name=file_name,
                            mime="application/octet-stream"
                        )
                    else:
                        st.error("Failed to generate summary. Please try again.")
                else:
                    st.error("No text could be extracted from the file. Please check the file format and content.")
    
    elif input_method == "✏️ Enter Text":
        st.markdown("### ✏️ Text Input")
        st.markdown("Type or paste your text for instant AI summarization")
        
        col1, col2 = st.columns(2)
        
        with col1:
            st.markdown("**Input Text**")
            input_text = st.text_area(
                "Enter your text here", 
                placeholder="Paste or type the text you want to summarize...",
                height=400,
                label_visibility="collapsed"
            )
        
        with col2:
            st.markdown("**Summary**")
            summary_placeholder = st.empty()
            
            if input_text.strip():
                with st.spinner("Generating summary..."):
                    summary = summarizer.generate_summary(input_text)
                summary_placeholder.text_area(
                    "Generated summary",
                    summary,
                    height=400,
                    disabled=True,
                    label_visibility="collapsed"
                )
            else:
                summary_placeholder.text_area(
                    "Generated summary",
                    "Enter text in the left box to see the summary here...",
                    height=400,
                    disabled=True,
                    label_visibility="collapsed"
                )
    
    elif input_method == "📧 Upload JSON (Email Threads)":
        st.markdown("### 📧 Email Thread Summarization")
        st.info("🎯 **Smart Email Processing**: Upload JSON files containing email thread data for intelligent conversation summaries.")
        
        uploaded_json = st.file_uploader(
            "Choose a JSON file", 
            type=["json"],
            help="JSON file should contain email thread data with thread_id and body fields"
        )
        
        if uploaded_json is not None:
            st.success(f"✅ JSON file '{uploaded_json.name}' uploaded successfully!")
            
            if st.button("🚀 Generate Thread Summaries", type="primary"):
                with st.spinner("Processing JSON file..."):
                    json_data = process_json_file(uploaded_json)
                
                if json_data:
                    with st.spinner("Generating summaries for email threads..."):
                        summaries = summarizer.summarize_json_threads(json_data)
                    
                    if summaries:
                        st.markdown("### 📝 Thread Summaries")
                        st.success(f"✨ **Generated {len(summaries)} thread summaries!**")
                        
                        # Display summaries
                        for i, summary in enumerate(summaries):
                            with st.expander(f"Thread {summary.get('thread_id', 'Unknown')}"):
                                st.write(summary.get('body', 'No summary available'))
                        
                        # Download functionality
                        json_content = json.dumps(summaries, indent=2)
                        
                        st.download_button(
                            label="⬇️ Download JSON Summary",
                            data=json_content.encode('utf-8'),
                            file_name="thread_summaries.json",
                            mime="application/json"
                        )
                    else:
                        st.error("Failed to generate summaries. Please check your JSON file format.")
                else:
                    st.error("Failed to process JSON file. Please check the file format.")
    
    # Enhanced Footer
    st.markdown("---")
    st.markdown("""
    <div style='text-align: center; padding: 2rem; background: linear-gradient(90deg, #667eea, #764ba2); 
                border-radius: 10px; margin-top: 2rem;'>
        <h3 style='color: white; margin-bottom: 1rem;'>🌟 Experience REEASE Everywhere</h3>
        <div style='display: flex; justify-content: center; gap: 2rem; margin-bottom: 1rem;'>
            <div style='color: white;'>
                <strong>🌐 Web App</strong><br>
                <small>Multi-format processing</small>
            </div>
            <div style='color: white;'>
                <strong>🔧 Chrome Extension</strong><br>
                <small>Browser integration</small>
            </div>
            <div style='color: white;'>
                <strong>🤖 AI Powered</strong><br>
                <small>Google Gemini 2.5</small>
            </div>
        </div>
        <p style='color: rgba(255,255,255,0.8); margin: 0; font-size: 0.9em;'>
            Built with ❤️ using Streamlit • Powered by Google Gemini AI
        </p>
    </div>
    """, unsafe_allow_html=True)

if __name__ == "__main__":
    main()
